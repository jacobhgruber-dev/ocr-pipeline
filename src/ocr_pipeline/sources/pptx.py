"""PowerPoint document source — extract text from slide decks.

Uses ``python-pptx`` to extract text from each slide.  Each slide is
treated as a logical page.  Rendering a slide to PNG is supported via
a placeholder approach.
"""

from __future__ import annotations

import logging
from pathlib import Path

from ocr_pipeline.errors import RenderError

from .base import DocumentSource

logger = logging.getLogger(__name__)


class PptxSource(DocumentSource):
    """Document source for PowerPoint files (.pptx).

    Each slide is a single logical page.  Text extraction walks all
    shapes on each slide and aggregates text content.  Rendering is
    not supported — convert slides to images before OCR.
    """

    _slides_text: list[str] | None = None

    @property
    def source_format(self) -> str:
        return "pptx"

    @property
    def source_mimetype(self) -> str:
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def _load_slides(self) -> list[str]:
        """Extract text from each slide, caching the result."""
        if self._slides_text is not None:
            return self._slides_text

        from pptx import Presentation

        try:
            prs = Presentation(str(self.path))
        except Exception as exc:
            raise RenderError(f"Failed to open PPTX: {self.path}") from exc
        slides: list[str] = []

        for slide_idx, slide in enumerate(prs.slides):
            parts: list[str] = [f"## Slide {slide_idx + 1}", ""]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
            # Extract speaker notes if present
            try:
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        parts.append("")
                        parts.append(f"**Notes:** {notes_text}")
            except Exception:
                pass
            slides.append("\n".join(parts))

        self._slides_text = slides
        return slides

    @property
    def page_count(self) -> int:
        try:
            return len(self._load_slides())
        except Exception:
            return 0

    def render_page(self, page_index: int, output_dir: Path, dpi: int = 300) -> Path:
        """Render a PPTX slide to PNG via LibreOffice → fitz.

        Converts to PDF using LibreOffice headless, then renders the slide.
        """
        import shutil
        import subprocess
        import tempfile

        soffice = shutil.which("soffice")
        if not soffice:
            raise NotImplementedError(
                "PptxSource.render_page requires LibreOffice (https://www.libreoffice.org). "
                "Install with: brew install libreoffice (macOS) or apt install libreoffice (Linux). "
                "Without it, use the text extraction path (no rendering needed for text-based PPTX)."
            )

        output_dir.mkdir(parents=True, exist_ok=True)
        png_path = output_dir / f"page_{page_index + 1:04d}.png"
        if png_path.exists():
            return png_path

        with tempfile.TemporaryDirectory() as td:
            pptx_tmp = Path(td) / self.path.name
            shutil.copy(str(self.path), str(pptx_tmp))
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", td, str(pptx_tmp)],
                capture_output=True,
                text=True,
                timeout=60,
            )
            pdf_path = pptx_tmp.with_suffix(".pdf")
            if result.returncode != 0 or not pdf_path.exists():
                raise RenderError(
                    f"LibreOffice conversion failed for {self.path.name}: {result.stderr[:200]}"
                )

            from ocr_pipeline.renderer import render_page as _render

            return _render(pdf_path, page_index, output_dir, dpi=dpi)

    def extract_text(
        self, page_index: int, output_dir: Path, flags: int | None = None
    ) -> tuple[str, Path | None]:
        slides = self._load_slides()
        if page_index < 0 or page_index >= len(slides):
            raise RenderError(f"PPTX slide index {page_index} out of range ({len(slides)} slides)")

        text = slides[page_index]

        output_dir.mkdir(parents=True, exist_ok=True)
        out_path = output_dir / f"page_{page_index + 1:04d}_final.md"
        out_path.write_text(text, encoding="utf-8")

        return text, out_path
