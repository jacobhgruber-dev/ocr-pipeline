"""Tests for multi-format document sources (DocumentSource implementations)."""

from __future__ import annotations

import json
import shutil
import zipfile
from pathlib import Path

import pytest

from ocr_pipeline.errors import ConfigError
from ocr_pipeline.sources import (
    ArchiveSource,
    ComicSource,
    CsvSource,
    DjvuSource,
    DocxSource,
    EmailSource,
    EpubSource,
    ExcelSource,
    HtmlSource,
    ImageSource,
    JsonSource,
    LatexSource,
    MarkdownSource,
    NotebookSource,
    OdtSource,
    PdfSource,
    PptxSource,
    RtfSource,
    SubtitleSource,
    TxtSource,
    detect_source,
)


# ---------------------------------------------------------------------------
# Factory detection
# ---------------------------------------------------------------------------


class TestDetectSource:
    def test_detect_pdf(self) -> None:
        pdf = Path("tests/fixtures/general/mixed_format.pdf")
        source = detect_source(pdf)
        assert isinstance(source, PdfSource)
        assert source.source_format == "pdf"
        assert source.page_count >= 1

    def test_detect_txt_by_extension(self, tmp_path: Path) -> None:
        f = tmp_path / "test.txt"
        f.write_text("hello world", encoding="utf-8")
        source = detect_source(f)
        assert isinstance(source, TxtSource)
        assert source.source_format == "txt"

    def test_detect_unsupported_raises(self, tmp_path: Path) -> None:
        f = tmp_path / "test.xyz"
        f.write_text("nope")
        with pytest.raises(ConfigError, match="Unsupported"):
            detect_source(f)


# ---------------------------------------------------------------------------
# PdfSource
# ---------------------------------------------------------------------------


class TestPdfSource:
    def test_basic(self) -> None:
        pdf = Path("tests/fixtures/general/mixed_format.pdf")
        source = PdfSource(pdf)
        assert source.source_format == "pdf"
        assert source.source_mimetype == "application/pdf"
        assert source.page_count >= 1

    def test_render_page(self, tmp_path: Path) -> None:
        pdf = Path("tests/fixtures/general/mixed_format.pdf")
        source = PdfSource(pdf)
        png = source.render_page(0, tmp_path, dpi=72)
        assert png.exists()
        assert png.suffix == ".png"
        assert png.stat().st_size > 100

    def test_extract_text(self, tmp_path: Path) -> None:
        pdf = Path("tests/fixtures/general/mixed_format.pdf")
        source = PdfSource(pdf)
        text, saved = source.extract_text(0, tmp_path)
        assert isinstance(text, str)
        assert saved is not None
        assert saved.exists()

    def test_extract_text_empty_page(self, tmp_path: Path) -> None:
        """Out-of-range page raises RenderError."""
        from ocr_pipeline.errors import RenderError

        pdf = Path("tests/fixtures/general/mixed_format.pdf")
        source = PdfSource(pdf)
        with pytest.raises(RenderError, match="out of range"):
            source.extract_text(999, tmp_path)


# ---------------------------------------------------------------------------
# TxtSource
# ---------------------------------------------------------------------------


class TestTxtSource:
    def test_basic_utf8(self, tmp_path: Path) -> None:
        f = tmp_path / "doc.txt"
        f.write_text("Hello world\nLine two\n", encoding="utf-8")
        source = TxtSource(f)
        assert source.source_format == "txt"
        assert source.page_count == 1

    def test_extract_text(self, tmp_path: Path) -> None:
        f = tmp_path / "doc.txt"
        f.write_text("test content", encoding="utf-8")
        source = TxtSource(f)
        out_dir = tmp_path / "out"
        text, saved = source.extract_text(0, out_dir)
        assert "test content" in text
        assert saved is not None
        assert saved.exists()

    def test_cached_read(self, tmp_path: Path) -> None:
        f = tmp_path / "doc.txt"
        f.write_text("cached", encoding="utf-8")
        source = TxtSource(f)
        text1, _ = source.extract_text(0, tmp_path / "out1")
        text2, _ = source.extract_text(0, tmp_path / "out2")
        assert text1 == text2 == "cached"

    def test_render_page_raises(self, tmp_path: Path) -> None:
        f = tmp_path / "doc.txt"
        f.write_text("text", encoding="utf-8")
        source = TxtSource(f)
        with pytest.raises(NotImplementedError, match="not supported"):
            source.render_page(0, tmp_path)


# ---------------------------------------------------------------------------
# ImageSource
# ---------------------------------------------------------------------------


class TestImageSource:
    def test_basic(self) -> None:
        """Use an existing render PNG from fixtures if available."""
        renders = list(Path("tests/fixtures").rglob("page_0001.png"))
        if not renders:
            pytest.skip("No test PNG found in fixtures")
        source = ImageSource(renders[0])
        assert source.source_format == "image"
        assert source.page_count == 1

    def test_extract_text_returns_empty(self) -> None:
        """Images have no extractable text — OCR must handle them."""
        renders = list(Path("tests/fixtures").rglob("page_0001.png"))
        if not renders:
            pytest.skip("No test PNG found in fixtures")
        source = ImageSource(renders[0])
        text, saved = source.extract_text(0, Path("/tmp"))
        assert text == ""
        assert saved is None

    def test_render_page_produces_output(self, tmp_path: Path) -> None:
        """ImageSource.render_page saves a PNG to the output directory."""
        renders = list(Path("tests/fixtures").rglob("page_0001.png"))
        if not renders:
            pytest.skip("No test PNG found in fixtures")
        source = ImageSource(renders[0])
        png = source.render_page(0, tmp_path)
        assert png.exists()
        assert png.suffix == ".png"
        assert png.stat().st_size > 100


# ---------------------------------------------------------------------------
# DocxSource
# ---------------------------------------------------------------------------


class TestDocxSource:
    @pytest.fixture()
    def docx_path(self, tmp_path: Path) -> Path:
        try:
            from docx import Document
        except ImportError:
            pytest.skip("python-docx not installed")

        doc = Document()
        doc.core_properties.title = "Test Title"
        doc.core_properties.author = "Test Author"
        doc.add_heading("Heading 1", level=1)
        doc.add_paragraph("This is a test paragraph.")
        doc.add_paragraph("Second paragraph with more text.")

        path = tmp_path / "test.docx"
        doc.save(str(path))
        return path

    def test_basic(self, docx_path: Path) -> None:
        source = DocxSource(docx_path)
        assert source.source_format == "docx"
        assert source.page_count == 1

    def test_extract_text(self, docx_path: Path, tmp_path: Path) -> None:
        source = DocxSource(docx_path)
        text, saved = source.extract_text(0, tmp_path)
        assert "Heading 1" in text
        assert "test paragraph" in text
        assert saved is not None


# ---------------------------------------------------------------------------
# CsvSource
# ---------------------------------------------------------------------------


class TestCsvSource:
    def test_basic(self, tmp_path: Path) -> None:
        f = tmp_path / "test.csv"
        f.write_text("name,age,city\nAlice,30,NYC\nBob,25,LA\n", encoding="utf-8")
        source = CsvSource(f)
        assert source.source_format == "csv"
        assert source.page_count == 1

    def test_extract_text(self, tmp_path: Path) -> None:
        f = tmp_path / "test.csv"
        f.write_text("col1,col2\na,b\nc,d\n", encoding="utf-8")
        source = CsvSource(f)
        text, saved = source.extract_text(0, tmp_path / "out")
        assert "col1" in text
        assert "col2" in text
        assert "a | b" in text.lower().replace(" ", "") or "a" in text
        assert saved is not None

    def test_extract_text_empty(self, tmp_path: Path) -> None:
        f = tmp_path / "empty.csv"
        f.write_text("col1,col2\n", encoding="utf-8")
        source = CsvSource(f)
        text, _ = source.extract_text(0, tmp_path / "out")
        assert isinstance(text, str)  # doesn't crash


# ---------------------------------------------------------------------------
# EpubSource
# ---------------------------------------------------------------------------


class TestEpubSource:
    @pytest.fixture()
    def epub_path(self, tmp_path: Path) -> Path:
        try:
            from ebooklib import epub
        except ImportError:
            pytest.skip("EbookLib not installed")

        book = epub.EpubBook()
        book.set_identifier("test-123")
        book.set_title("Test Book")
        book.set_language("en")
        book.add_author("Jane Author")

        ch1 = epub.EpubHtml(title="Chapter 1", file_name="ch1.xhtml")
        ch1.content = (
            "<html><body><h1>Chapter One</h1><p>This is chapter one text.</p></body></html>"
        )
        book.add_item(ch1)

        ch2 = epub.EpubHtml(title="Chapter 2", file_name="ch2.xhtml")
        ch2.content = "<html><body><h1>Chapter Two</h1><p>Chapter two content.</p></body></html>"
        book.add_item(ch2)

        book.spine = ["nav", ch1, ch2]
        book.toc = [ch1, ch2]
        book.add_item(epub.EpubNcx())
        book.add_item(epub.EpubNav())

        path = tmp_path / "test.epub"
        epub.write_epub(str(path), book)
        return path

    def test_basic(self, epub_path: Path) -> None:
        source = EpubSource(epub_path)
        assert source.source_format == "epub"
        assert source.page_count >= 1

    def test_extract_text(self, epub_path: Path, tmp_path: Path) -> None:
        source = EpubSource(epub_path)
        text, saved = source.extract_text(0, tmp_path)
        assert isinstance(text, str)
        assert len(text) > 0
        assert saved is not None


# ---------------------------------------------------------------------------
# PptxSource
# ---------------------------------------------------------------------------


class TestPptxSource:
    @pytest.fixture()
    def pptx_path(self, tmp_path: Path) -> Path:
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        prs = Presentation()
        prs.core_properties.title = "Test Presentation"

        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "Slide 1 Title"
        slide1.notes_slide.notes_text_frame.text = "Speaker notes for slide 1"

        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Slide 2 Title"

        path = tmp_path / "test.pptx"
        prs.save(str(path))
        return path

    def test_basic(self, pptx_path: Path) -> None:
        source = PptxSource(pptx_path)
        assert source.source_format == "pptx"
        assert source.page_count == 2  # 2 slides

    def test_extract_text(self, pptx_path: Path, tmp_path: Path) -> None:
        source = PptxSource(pptx_path)
        text, saved = source.extract_text(0, tmp_path)
        assert "Slide 1 Title" in text
        assert "Speaker notes" in text
        assert saved is not None

    def test_extract_text_page2(self, pptx_path: Path, tmp_path: Path) -> None:
        source = PptxSource(pptx_path)
        text, saved = source.extract_text(1, tmp_path / "p2")
        assert "Slide 2 Title" in text
        assert saved is not None


# ---------------------------------------------------------------------------
# ExcelSource
# ---------------------------------------------------------------------------


class TestExcelSource:
    @pytest.fixture()
    def xlsx_path(self, tmp_path: Path) -> Path:
        try:
            from openpyxl import Workbook
        except ImportError:
            pytest.skip("openpyxl not installed")

        wb = Workbook()
        wb.properties.title = "Test Workbook"
        wb.properties.creator = "Test Creator"

        ws = wb.active
        ws.title = "Data"
        ws.append(["Name", "Value", "Date"])
        ws.append(["Alpha", "10", "2024-01-01"])
        ws.append(["Beta", "20", "2024-06-15"])

        ws2 = wb.create_sheet("Summary")
        ws2.append(["Total", "30"])

        path = tmp_path / "test.xlsx"
        wb.save(str(path))
        return path

    def test_basic(self, xlsx_path: Path) -> None:
        source = ExcelSource(xlsx_path)
        assert source.source_format == "excel"
        assert source.page_count == 2  # 2 sheets

    def test_extract_text(self, xlsx_path: Path, tmp_path: Path) -> None:
        source = ExcelSource(xlsx_path)
        text, saved = source.extract_text(0, tmp_path)
        assert "Alpha" in text
        assert "Value" in text
        assert saved is not None

    def test_extract_text_sheet2(self, xlsx_path: Path, tmp_path: Path) -> None:
        source = ExcelSource(xlsx_path)
        text, saved = source.extract_text(1, tmp_path / "s2")
        assert "Total" in text
        assert saved is not None


# ---------------------------------------------------------------------------
# MarkdownSource
# ---------------------------------------------------------------------------


class TestMarkdownSource:
    def test_basic_no_frontmatter(self, tmp_path: Path) -> None:
        f = tmp_path / "test.md"
        f.write_text("# Hello\n\nThis is markdown.", encoding="utf-8")
        source = MarkdownSource(f)
        assert source.source_format == "markdown"
        assert source.page_count == 1

    def test_extract_text(self, tmp_path: Path) -> None:
        f = tmp_path / "test.md"
        f.write_text("# Title\n\nBody text here.", encoding="utf-8")
        source = MarkdownSource(f)
        text, saved = source.extract_text(0, tmp_path / "out")
        assert "Title" in text
        assert "Body text" in text
        assert saved is not None

    def test_frontmatter_parsing(self, tmp_path: Path) -> None:
        f = tmp_path / "test.md"
        f.write_text(
            "---\ntitle: My Document\nauthor: Jane Doe\ndate: 2025-01-15\n---\n\n# Content\n\nText here.",
            encoding="utf-8",
        )
        source = MarkdownSource(f)
        meta = source.extract_metadata()
        assert meta.title == "My Document"
        assert meta.authors == ["Jane Doe"]
        assert meta.date == "2025-01-15"
        assert meta.extraction_method == "frontmatter"

    def test_html_comments_stripped(self, tmp_path: Path) -> None:
        f = tmp_path / "test.md"
        f.write_text("<!-- doc: title: Test | author: Me -->\n\n# Actual content", encoding="utf-8")
        source = MarkdownSource(f)
        text, _ = source.extract_text(0, tmp_path / "out")
        assert "Actual content" in text
        assert "doc:" not in text  # HTML comment stripped


# ---------------------------------------------------------------------------
# HtmlSource
# ---------------------------------------------------------------------------


class TestHtmlSource:
    def test_basic(self, tmp_path: Path) -> None:
        f = tmp_path / "test.html"
        f.write_text(
            "<html><head><title>Test Page</title></head><body><h1>Heading</h1><p>Paragraph.</p></body></html>",
            encoding="utf-8",
        )
        source = HtmlSource(f)
        assert source.source_format == "html"
        assert source.page_count == 1

    def test_extract_text(self, tmp_path: Path) -> None:
        f = tmp_path / "test.html"
        f.write_text("<html><body><p>Hello world</p></body></html>", encoding="utf-8")
        source = HtmlSource(f)
        text, saved = source.extract_text(0, tmp_path / "out")
        assert "Hello world" in text
        assert saved is not None

    def test_json_ld_metadata(self, tmp_path: Path) -> None:
        f = tmp_path / "test.html"
        f.write_text(
            """<html><head>
<script type="application/ld+json">
{"@context": "https://schema.org", "@type": "ScholarlyArticle",
 "name": "Test Article", "author": {"name": "Dr. Smith"},
 "datePublished": "2024-06-01", "description": "A test paper."}
</script></head><body>Content</body></html>""",
            encoding="utf-8",
        )
        source = HtmlSource(f)
        meta = source.extract_metadata()
        assert meta.title == "Test Article"
        assert meta.authors == ["Dr. Smith"]
        assert meta.date == "2024-06-01"
        assert meta.extraction_method == "html-metadata"

    def test_meta_tags(self, tmp_path: Path) -> None:
        f = tmp_path / "test.html"
        f.write_text(
            "<html><head>"
            '<meta name="citation_title" content="Cited Paper">'
            '<meta name="citation_author" content="A. Author">'
            '<meta name="dc.publisher" content="University Press">'
            "</head><body>Text</body></html>",
            encoding="utf-8",
        )
        source = HtmlSource(f)
        meta = source.extract_metadata()
        assert meta.title == "Cited Paper"


# ---------------------------------------------------------------------------
# LatexSource
# ---------------------------------------------------------------------------


class TestLatexSource:
    def test_basic(self, tmp_path: Path) -> None:
        f = tmp_path / "test.tex"
        f.write_text(
            r"\title{My Paper}\author{A. U. Thor}\date{2025}\begin{document}\section{Intro}Hello world.\end{document}",
            encoding="utf-8",
        )
        source = LatexSource(f)
        assert source.source_format == "latex"
        assert source.page_count == 1

    def test_extract_text_strips_commands(self, tmp_path: Path) -> None:
        f = tmp_path / "test.tex"
        f.write_text(
            r"\documentclass{article}\begin{document}\textbf{Bold} and \emph{italic} text.\end{document}",
            encoding="utf-8",
        )
        source = LatexSource(f)
        text, saved = source.extract_text(0, tmp_path / "out")
        assert "Bold" in text
        # Strip commands should remove \textbf and \emph but keep the content
        assert "italic" in text.lower()
        assert saved is not None

    def test_metadata_extraction(self, tmp_path: Path) -> None:
        f = tmp_path / "test.tex"
        f.write_text(
            r"\title{Quantum Computing Primer}"
            r"\author{Alice Researcher, Bob Scientist}"
            r"\date{March 2025}"
            r"\begin{abstract}This paper introduces quantum computing.\end{abstract}"
            r"\begin{document}\section{Body}Content.\end{document}",
            encoding="utf-8",
        )
        source = LatexSource(f)
        meta = source.extract_metadata()
        assert "Quantum Computing" in meta.title
        assert len(meta.authors) == 1
        assert meta.extraction_method == "latex-parsing"

    def test_comments_stripped(self, tmp_path: Path) -> None:
        f = tmp_path / "test.tex"
        f.write_text(
            "% This is a comment\n\\begin{document}Visible content\\end{document}",
            encoding="utf-8",
        )
        source = LatexSource(f)
        text, _ = source.extract_text(0, tmp_path / "out")
        assert "Visible content" in text
        assert "comment" not in text.lower()

    def test_render_page_raises(self, tmp_path: Path) -> None:
        f = tmp_path / "test.tex"
        f.write_text(r"\begin{document}text\end{document}", encoding="utf-8")
        source = LatexSource(f)
        with pytest.raises(NotImplementedError):
            source.render_page(0, tmp_path)


# ---------------------------------------------------------------------------
# JsonSource
# ---------------------------------------------------------------------------


class TestJsonSource:
    def test_basic(self, tmp_path: Path) -> None:
        f = tmp_path / "test.json"
        f.write_text('{"key": "value"}', encoding="utf-8")
        source = JsonSource(f)
        assert source.source_format == "json"
        assert source.source_mimetype == "application/json"
        assert source.page_count == 1

    def test_extract_text(self, tmp_path: Path) -> None:
        f = tmp_path / "test.json"
        f.write_text('{"name": "test", "items": [1, 2, 3]}', encoding="utf-8")
        source = JsonSource(f)
        text, saved = source.extract_text(0, tmp_path / "out")
        assert "test" in text
        assert "items" in text
        assert saved is not None
        assert saved.exists()

    def test_json_ld_metadata(self, tmp_path: Path) -> None:
        f = tmp_path / "test.json"
        f.write_text(
            json.dumps(
                {
                    "@context": "https://schema.org",
                    "@type": "ScholarlyArticle",
                    "name": "Test Article",
                    "author": {"name": "Dr. Smith"},
                    "datePublished": "2024-06-01",
                    "description": "A test paper.",
                }
            ),
            encoding="utf-8",
        )
        source = JsonSource(f)
        meta = source.extract_metadata()
        assert meta.title == "Test Article"
        assert meta.authors == ["Dr. Smith"]
        assert meta.date == "2024-06-01"
        assert meta.extraction_method == "json-parsing"

    def test_pure_json_no_json_ld(self, tmp_path: Path) -> None:
        """Pure JSON (no JSON-LD) still extracts text successfully."""
        f = tmp_path / "test.json"
        f.write_text('{"data": {"value": 42}}', encoding="utf-8")
        source = JsonSource(f)
        text, _ = source.extract_text(0, tmp_path / "out")
        assert "42" in text
        assert "value" in text
        # Metadata should still work without JSON-LD
        meta = source.extract_metadata()
        assert meta.extraction_method == "json-parsing"

    def test_render_page_raises(self, tmp_path: Path) -> None:
        f = tmp_path / "test.json"
        f.write_text("{}", encoding="utf-8")
        source = JsonSource(f)
        with pytest.raises(NotImplementedError):
            source.render_page(0, tmp_path)


# ---------------------------------------------------------------------------
# RtfSource
# ---------------------------------------------------------------------------


class TestRtfSource:
    def test_basic(self, tmp_path: Path) -> None:
        f = tmp_path / "test.rtf"
        f.write_text("{\\rtf1\\ansi Hello World}", encoding="utf-8")
        source = RtfSource(f)
        assert source.source_format == "rtf"
        assert source.source_mimetype == "application/rtf"
        assert source.page_count == 1

    def test_extract_text(self, tmp_path: Path) -> None:
        """Text extraction strips RTF control codes."""
        f = tmp_path / "test.rtf"
        f.write_text("{\\rtf1\\ansi\\deff0 {\\b Bold} and normal text.}", encoding="utf-8")
        source = RtfSource(f)
        text, saved = source.extract_text(0, tmp_path / "out")
        assert "Bold" in text
        assert "normal text" in text
        assert saved is not None
        assert saved.exists()

    def test_metadata(self, tmp_path: Path) -> None:
        f = tmp_path / "test.rtf"
        f.write_text("{\\rtf1\\ansi minimal}", encoding="utf-8")
        source = RtfSource(f)
        meta = source.extract_metadata()
        assert meta.extraction_method == "rtf-stripping"
        assert meta.document_type == "document"

    def test_render_page_raises(self, tmp_path: Path) -> None:
        f = tmp_path / "test.rtf"
        f.write_text("{\\rtf1\\ansi test}", encoding="utf-8")
        source = RtfSource(f)
        with pytest.raises(NotImplementedError):
            source.render_page(0, tmp_path)


# ---------------------------------------------------------------------------
# OdtSource
# ---------------------------------------------------------------------------


class TestOdtSource:
    """ODT files are ZIP containers with content.xml and meta.xml."""

    @pytest.fixture()
    def odt_path(self, tmp_path: Path) -> Path:
        """Create a minimal valid ODT file using zipfile."""
        path = tmp_path / "test.odt"
        with zipfile.ZipFile(str(path), "w") as zf:
            # mimetype must be first and stored uncompressed for valid ODT
            zf.writestr("mimetype", "application/vnd.oasis.opendocument.text")
            zf.writestr(
                "content.xml",
                """<?xml version="1.0" encoding="UTF-8"?>
<office:document-content
 xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"
 xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">
  <office:body>
    <office:text>
      <text:p>Hello ODT World</text:p>
      <text:h>ODT Heading</text:h>
    </office:text>
  </office:body>
</office:document-content>""",
            )
            zf.writestr(
                "meta.xml",
                """<?xml version="1.0" encoding="UTF-8"?>
<office:document-meta
 xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"
 xmlns:meta="urn:oasis:names:tc:opendocument:xmlns:meta:1.0"
 xmlns:dc="http://purl.org/dc/elements/1.1/">
  <office:meta>
    <dc:title>Test ODT Document</dc:title>
    <dc:creator>ODT Author</dc:creator>
    <meta:document-statistic meta:page-count="1"/>
  </office:meta>
</office:document-meta>""",
            )
        return path

    def test_basic(self, odt_path: Path) -> None:
        source = OdtSource(odt_path)
        assert source.source_format == "odt"
        assert source.page_count == 1

    def test_extract_text(self, odt_path: Path, tmp_path: Path) -> None:
        source = OdtSource(odt_path)
        text, saved = source.extract_text(0, tmp_path)
        assert "Hello ODT World" in text
        assert "ODT Heading" in text
        assert saved is not None
        assert saved.exists()

    def test_metadata(self, odt_path: Path) -> None:
        source = OdtSource(odt_path)
        meta = source.extract_metadata()
        assert meta.title == "Test ODT Document"
        assert meta.authors == ["ODT Author"]
        assert meta.extraction_method == "odt-parsing"

    def test_render_page_raises(self, odt_path: Path, tmp_path: Path) -> None:
        source = OdtSource(odt_path)
        with pytest.raises(NotImplementedError):
            source.render_page(0, tmp_path)


# ---------------------------------------------------------------------------
# NotebookSource
# ---------------------------------------------------------------------------


class TestNotebookSource:
    @pytest.fixture()
    def notebook_path(self, tmp_path: Path) -> Path:
        """Create a minimal .ipynb with markdown + code cells."""
        nb = {
            "metadata": {
                "kernelspec": {"display_name": "Python 3", "name": "python3"},
                "language_info": {"name": "python"},
                "title": "My Notebook",
                "authors": ["Alice Coder", "Bob Analyst"],
            },
            "nbformat": 4,
            "nbformat_minor": 5,
            "cells": [
                {
                    "cell_type": "markdown",
                    "source": ["# Introduction\n", "This is a notebook about data."],
                },
                {
                    "cell_type": "code",
                    "source": ["import pandas as pd\n", "df = pd.read_csv('data.csv')"],
                    "outputs": [
                        {
                            "output_type": "stream",
                            "text": ["Hello from stdout\n"],
                        }
                    ],
                },
            ],
        }
        path = tmp_path / "test.ipynb"
        path.write_text(json.dumps(nb), encoding="utf-8")
        return path

    def test_basic(self, notebook_path: Path) -> None:
        source = NotebookSource(notebook_path)
        assert source.source_format == "notebook"
        assert source.page_count == 1

    def test_extract_text(self, notebook_path: Path, tmp_path: Path) -> None:
        source = NotebookSource(notebook_path)
        text, saved = source.extract_text(0, tmp_path)
        assert "Introduction" in text
        assert "import pandas" in text
        assert "Hello from stdout" in text
        assert saved is not None

    def test_metadata(self, notebook_path: Path) -> None:
        source = NotebookSource(notebook_path)
        meta = source.extract_metadata()
        assert meta.title == "My Notebook"
        assert meta.authors == ["Alice Coder, Bob Analyst"]
        assert meta.language == "python"
        assert meta.extraction_method == "notebook-parsing"

    def test_render_page_raises(self, notebook_path: Path, tmp_path: Path) -> None:
        source = NotebookSource(notebook_path)
        with pytest.raises(NotImplementedError):
            source.render_page(0, tmp_path)


# ---------------------------------------------------------------------------
# ArchiveSource
# ---------------------------------------------------------------------------


class TestArchiveSource:
    @pytest.fixture()
    def zip_path(self, tmp_path: Path) -> Path:
        """Create a minimal .zip containing text files and a readme."""
        path = tmp_path / "test.zip"
        with zipfile.ZipFile(str(path), "w") as zf:
            zf.writestr("readme.md", "# Test Archive\n\nThis is a test readme.")
            zf.writestr("data.txt", "sample data content")
            zf.writestr("notes/note.txt", "note content")
        return path

    def test_basic(self, zip_path: Path) -> None:
        source = ArchiveSource(zip_path)
        assert source.source_format == "zip"
        assert source.source_mimetype == "application/zip"
        assert source.page_count == 1

    def test_extract_text_lists_files(self, zip_path: Path, tmp_path: Path) -> None:
        source = ArchiveSource(zip_path)
        text, saved = source.extract_text(0, tmp_path)
        assert "readme.md" in text
        assert "data.txt" in text
        assert "Test Archive" in text  # readme content included
        assert saved is not None

    def test_metadata(self, zip_path: Path) -> None:
        source = ArchiveSource(zip_path)
        meta = source.extract_metadata()
        assert meta.extraction_method == "archive-listing"
        assert meta.document_type == "archive"
        assert meta.source_info.extra.get("file_count") == 3

    def test_render_page_raises(self, zip_path: Path, tmp_path: Path) -> None:
        source = ArchiveSource(zip_path)
        with pytest.raises(NotImplementedError):
            source.render_page(0, tmp_path)


# ---------------------------------------------------------------------------
# EmailSource
# ---------------------------------------------------------------------------


class TestEmailSource:
    @pytest.fixture()
    def eml_path(self, tmp_path: Path) -> Path:
        """Create a minimal .eml file with headers and body."""
        path = tmp_path / "test.eml"
        path.write_text(
            "From: sender@example.com\r\n"
            "To: recipient@example.com\r\n"
            "Subject: Test Email\r\n"
            "Date: Thu, 01 Jan 2025 00:00:00 +0000\r\n"
            "Content-Type: text/plain; charset=utf-8\r\n"
            "\r\n"
            "This is the email body.\r\n",
            encoding="utf-8",
        )
        return path

    def test_basic(self, eml_path: Path) -> None:
        source = EmailSource(eml_path)
        assert source.source_format == "email"
        assert source.page_count == 1

    def test_extract_text(self, eml_path: Path, tmp_path: Path) -> None:
        source = EmailSource(eml_path)
        text, saved = source.extract_text(0, tmp_path)
        assert "This is the email body." in text
        assert "Subject: Test Email" in text
        assert "From: sender@example.com" in text
        assert saved is not None

    def test_metadata(self, eml_path: Path) -> None:
        source = EmailSource(eml_path)
        meta = source.extract_metadata()
        assert meta.title == "Test Email"
        assert meta.authors == ["sender@example.com"]
        assert meta.extraction_method == "email-parsing"
        assert meta.document_type == "email"

    def test_render_page_raises(self, eml_path: Path, tmp_path: Path) -> None:
        source = EmailSource(eml_path)
        with pytest.raises(NotImplementedError):
            source.render_page(0, tmp_path)


# ---------------------------------------------------------------------------
# SubtitleSource
# ---------------------------------------------------------------------------


class TestSubtitleSource:
    def test_basic(self, tmp_path: Path) -> None:
        f = tmp_path / "test.srt"
        f.write_text(
            "1\n"
            "00:00:01,000 --> 00:00:04,000\n"
            "Hello world\n"
            "\n"
            "2\n"
            "00:00:05,000 --> 00:00:08,000\n"
            "This is a test\n",
            encoding="utf-8",
        )
        source = SubtitleSource(f)
        assert source.source_format == "srt"
        assert source.source_mimetype == "application/x-subrip"
        assert source.page_count == 1

    def test_extract_text(self, tmp_path: Path) -> None:
        f = tmp_path / "test.srt"
        f.write_text(
            "1\n"
            "00:00:01,000 --> 00:00:04,000\n"
            "Hello world\n"
            "\n"
            "2\n"
            "00:00:05,000 --> 00:00:08,000\n"
            "This is a test\n",
            encoding="utf-8",
        )
        source = SubtitleSource(f)
        text, saved = source.extract_text(0, tmp_path / "out")
        assert "Hello world" in text
        assert "This is a test" in text
        # Timestamps and sequence numbers are stripped
        assert "00:00:01,000" not in text
        assert "1" not in text.splitlines()
        assert saved is not None

    def test_metadata(self, tmp_path: Path) -> None:
        f = tmp_path / "test.srt"
        f.write_text(
            "1\n00:00:01,000 --> 00:00:04,000\nHello world\n",
            encoding="utf-8",
        )
        source = SubtitleSource(f)
        meta = source.extract_metadata()
        assert meta.extraction_method == "subtitle-parsing"
        assert meta.document_type == "subtitle"

    def test_render_page_raises(self, tmp_path: Path) -> None:
        f = tmp_path / "test.srt"
        f.write_text(
            "1\n00:00:01,000 --> 00:00:04,000\nTest\n",
            encoding="utf-8",
        )
        source = SubtitleSource(f)
        with pytest.raises(NotImplementedError):
            source.render_page(0, tmp_path)


# ---------------------------------------------------------------------------
# DjvuSource
# ---------------------------------------------------------------------------


class TestDjvuSource:
    """Tests for DJVU document source (conditionally requires djvulibre CLI)."""

    def test_source_format_always(self, tmp_path: Path) -> None:
        """source_format returns 'djvu' even without djvulibre installed."""
        f = tmp_path / "test.djvu"
        f.write_text("")
        source = DjvuSource(f)
        assert source.source_format == "djvu"
        assert source.source_mimetype == "image/vnd.djvu"

    def test_detect_via_extension(self, tmp_path: Path) -> None:
        """detect_source routes .djvu and .djv files to DjvuSource."""
        for ext in (".djvu", ".djv"):
            f = tmp_path / f"test{ext}"
            f.write_text("")
            source = detect_source(f)
            assert isinstance(source, DjvuSource)
            assert source.source_format == "djvu"

    def test_page_count_fallback(self, tmp_path: Path) -> None:
        """Without djvudump, page_count falls back to 1."""
        f = tmp_path / "test.djvu"
        f.write_text("")
        source = DjvuSource(f)
        assert source.page_count >= 1

    def test_extract_text_basic(self, tmp_path: Path) -> None:
        """Extract text returns a result (either real text or a fallback message)."""
        f = tmp_path / "test.djvu"
        f.write_text("")
        source = DjvuSource(f)
        text, saved = source.extract_text(0, tmp_path / "out")
        assert isinstance(text, str)
        assert len(text) >= 0
        if saved is not None:
            assert saved.exists()

    def test_metadata(self, tmp_path: Path) -> None:
        f = tmp_path / "test.djvu"
        f.write_text("")
        source = DjvuSource(f)
        meta = source.extract_metadata()
        assert meta.document_type == "book"
        assert meta.extraction_method == "djvu-cli"
        assert meta.source_info is not None
        assert meta.source_info.format == "djvu"

    def test_basic_with_djvutxt(self, tmp_path: Path) -> None:
        """If djvutxt is available, test basic functionality with a dummy file.

        Note: A valid DJVU file is needed for meaningful text extraction,
        but format detection and page count work on any file.
        """
        if not shutil.which("djvutxt"):
            pytest.skip("djvutxt not installed")
        f = tmp_path / "test.djvu"
        f.write_text("")
        source = DjvuSource(f)
        assert source.source_format == "djvu"
        assert source.page_count >= 1

    def test_extract_text_out_of_range(self, tmp_path: Path) -> None:
        """Out-of-range page returns empty string."""
        f = tmp_path / "test.djvu"
        f.write_text("")
        source = DjvuSource(f)
        text, saved = source.extract_text(999, tmp_path / "out")
        assert text == ""
        assert saved is None

    def test_render_page_raises_without_ddjvu(self, tmp_path: Path) -> None:
        """render_page raises NotImplementedError when ddjvu is not installed."""
        if shutil.which("ddjvu"):
            pytest.skip("ddjvu is installed — skipping error-path test")
        f = tmp_path / "test.djvu"
        f.write_text("")
        source = DjvuSource(f)
        with pytest.raises(NotImplementedError, match="ddjvu not found"):
            source.render_page(0, tmp_path)


# ---------------------------------------------------------------------------
# ComicSource
# ---------------------------------------------------------------------------


class TestComicSource:
    """Tests for comic book archive source (CBZ/CBR)."""

    @staticmethod
    def _make_1px_png() -> bytes:
        """Create a minimal 1x1 pixel PNG in memory using Pillow."""
        from PIL import Image
        import io

        img = Image.new("RGB", (1, 1), color=(255, 0, 0))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()

    @pytest.fixture()
    def cbz_path(self, tmp_path: Path) -> Path:
        """Create a minimal .cbz file containing a single 1-pixel PNG."""
        path = tmp_path / "comic.cbz"
        png_data = self._make_1px_png()
        with zipfile.ZipFile(str(path), "w") as zf:
            zf.writestr("page_01.png", png_data)
        return path

    def test_basic(self, cbz_path: Path) -> None:
        source = ComicSource(cbz_path)
        assert source.source_format == "cbz"
        assert source.source_mimetype == "application/x-cbz"
        assert source.page_count == 1

    def test_extract_text_returns_image_description(self, cbz_path: Path, tmp_path: Path) -> None:
        source = ComicSource(cbz_path)
        text, saved = source.extract_text(0, tmp_path)
        assert isinstance(text, str)
        assert "Comic page 1" in text
        assert "page_01.png" in text
        assert saved is None  # comics have no text file to save

    def test_extract_text_out_of_range(self, cbz_path: Path, tmp_path: Path) -> None:
        source = ComicSource(cbz_path)
        text, saved = source.extract_text(999, tmp_path)
        assert text == ""
        assert saved is None

    def test_metadata(self, cbz_path: Path) -> None:
        source = ComicSource(cbz_path)
        meta = source.extract_metadata()
        assert meta.document_type == "comic"
        assert meta.extraction_method == "comic-archive"
        assert meta.source_info is not None
        assert meta.source_info.format == "cbz"
        assert meta.source_info.extra.get("image_count") == 1

    def test_render_page(self, cbz_path: Path, tmp_path: Path) -> None:
        source = ComicSource(cbz_path)
        png = source.render_page(0, tmp_path)
        assert png.exists()
        assert png.suffix == ".png"
        assert png.stat().st_size > 0

    def test_render_page_out_of_range(self, cbz_path: Path, tmp_path: Path) -> None:
        source = ComicSource(cbz_path)
        with pytest.raises(IndexError, match="out of range"):
            source.render_page(999, tmp_path)

    def test_cbr_format_detection(self, tmp_path: Path) -> None:
        """Format detection for .cbr (RAR archives — can't easily create, test format only)."""
        f = tmp_path / "comic.cbr"
        f.write_text("not a real rar")
        source = ComicSource(f)
        assert source.source_format == "cbr"
        assert source.source_mimetype == "application/x-cbr"
        # CBR with no actual RAR content: page_count falls back to 1
        assert source.page_count >= 1

    def test_detect_via_extension(self, cbz_path: Path) -> None:
        """detect_source routes .cbz files to ComicSource."""
        source = detect_source(cbz_path)
        assert isinstance(source, ComicSource)
        assert source.source_format == "cbz"
        assert source.page_count == 1
